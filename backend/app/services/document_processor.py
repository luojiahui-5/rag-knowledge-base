"""
文档解析 + 切片 + 向量化管道
"""
import os
from ..core.config import get_settings

settings = get_settings()


def parse_file(file_path: str, file_type: str) -> str:
    """根据文件类型自动选择解析器"""
    parsers = {
        "pdf": parse_pdf,
        "docx": parse_docx,
        "doc": parse_docx,
        "pptx": parse_pptx,
        "ppt": parse_pptx,
        "xlsx": parse_xlsx,
        "xls": parse_xlsx,
        "md": parse_text,
        "txt": parse_text,
        "html": parse_text,
        "htm": parse_text,
    }
    parser = parsers.get(file_type, parse_text)
    try:
        return parser(file_path)
    except Exception as e:
        return f"[解析失败: {str(e)}]"


def parse_pdf(file_path: str) -> str:
    """PDF 解析"""
    try:
        import fitz
        doc = fitz.open(file_path)
        parts = []
        for i, page in enumerate(doc):
            text = page.get_text().strip()
            if text:
                parts.append(f"--- 第 {i+1} 页 ---\n{text}")
        doc.close()
        return "\n\n".join(parts) if parts else "[PDF 文件无可提取的文字内容]"
    except ImportError:
        return "[PDF 解析需要安装 PyMuPDF: pip install pymupdf]"
    except Exception as e:
        return f"[PDF 解析失败: {e}]"


def parse_docx(file_path: str) -> str:
    """Word 文档解析"""
    try:
        from docx import Document as DocxDoc
        doc = DocxDoc(file_path)
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                # 检测标题样式
                if para.style and para.style.name and 'Heading' in para.style.name:
                    parts.append(f"\n## {para.text}\n")
                else:
                    parts.append(para.text)

        # 表格
        for i, table in enumerate(doc.tables):
            parts.append(f"\n--- 表格 {i+1} ---")
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                parts.append(" | ".join(cells))

        return "\n\n".join(parts) if parts else "[Word 文件无可提取的文字内容]"
    except ImportError:
        return "[Word 解析需要安装 python-docx: pip install python-docx]"
    except Exception as e:
        return f"[Word 解析失败: {e}]"


def parse_pptx(file_path: str) -> str:
    """PPT 解析"""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        slide_texts.append(" | ".join(cells))
            if slide_texts:
                parts.append(f"--- 幻灯片 {i} ---\n" + "\n".join(slide_texts))
        return "\n\n".join(parts) if parts else "[PPT 文件无可提取的文字内容]"
    except ImportError:
        return "[PPT 解析需要安装 python-pptx: pip install python-pptx]"
    except Exception as e:
        return f"[PPT 解析失败: {e}]"


def parse_xlsx(file_path: str) -> str:
    """Excel 解析"""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"--- 工作表: {sheet_name} ---")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    parts.append(" | ".join(cells))
        return "\n".join(parts) if parts else "[Excel 文件无数据]"
    except Exception as e:
        return f"[Excel 解析失败: {e}]"


def parse_text(file_path: str) -> str:
    """纯文本 / Markdown / HTML 直接读取"""
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception as e:
        return f"[文件读取失败: {e}]"


def chunk_text(text: str, chunk_size: int = None, overlap: int = None) -> list[dict]:
    """文本切片"""
    chunk_size = chunk_size or settings.CHUNK_SIZE
    overlap = overlap or settings.CHUNK_OVERLAP

    try:
        from langchain_text_splitters import RecursiveCharacterTextSplitter
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=overlap,
            separators=["\n\n", "\n", "。", ". ", "；", "; ", "，", ", ", " ", ""],
        )
        return [{"content": c, "metadata": {}} for c in splitter.split_text(text)]
    except ImportError:
        chunks = []
        for i in range(0, len(text), chunk_size - overlap):
            chunk = text[i:i + chunk_size]
            if len(chunk) > 20:
                chunks.append({"content": chunk, "metadata": {}})
        return chunks
